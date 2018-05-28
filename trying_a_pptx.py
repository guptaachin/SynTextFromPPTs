# this library can only open pptx file i.e > 2007 pptx files

from pptx import Presentation
import os

def main():

    lang_folder = preprocessing()

    print(lang_folder)
    folder_for_ppt = lang_folder
    for each_ppt in os.listdir(lang_folder):
        # each_ppt not in con_set and
        if  (each_ppt.endswith('ppt') or each_ppt.endswith('pptx')):
            current_ppt_path = os.path.join(folder_for_ppt, "text.pptx")

            f = open(current_ppt_path, "rb")

            try:
                prs = Presentation(f)
                print(len(prs.slides))
            except Exception as e:
                print('Exception raised = ',e)

            # for each_slide in prs.slides:
            #     print(each_slide)

def preprocessing():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("language", help="lang_ja,lang_ko,lang_es")
    args = parser.parse_args()
    CURR_LANG = args.language

    data_folder = os.path.join(os.getcwd(), 'data')
    lang_folder = os.path.join(data_folder, CURR_LANG)

    images_folder = os.path.join(lang_folder, 'images') # data/lang_ja/images
    
    # remove only for texting
    [os.remove(images_folder+'/'+r) for r in os.listdir(images_folder)]   
    ppt_folder = os.path.join(lang_folder, 'ppts')
    lang_folder = os.path.join(data_folder, CURR_LANG)
    return lang_folder


if __name__ == "__main__":
    main()